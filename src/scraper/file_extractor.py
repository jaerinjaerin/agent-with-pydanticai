"""
첨부파일 텍스트/이미지/표 추출기.

PDF, DOCX, XLSX, HWP, PPTX, MD, TXT 파일에서 텍스트, 표(마크다운), 이미지를 추출한다.
"""

from dataclasses import dataclass, field
from pathlib import Path


@dataclass
class ExtractedImage:
    """추출된 이미지 데이터."""
    data: bytes
    filename: str
    page_or_slide: int  # 1-based
    context: str  # 주변 텍스트 (섹션 제목 등)


@dataclass
class ExtractionResult:
    """파일에서 추출된 전체 결과 (텍스트 + 이미지)."""
    text: str  # 전체 텍스트 (표 마크다운 포함)
    images: list[ExtractedImage] = field(default_factory=list)


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".hwp", ".pptx", ".md", ".txt"}


def extract_text(file_path: Path) -> str:
    """파일 확장자에 따라 적절한 방법으로 텍스트를 추출한다."""
    ext = file_path.suffix.lower()

    if ext == ".pdf":
        return _extract_pdf(file_path)
    elif ext == ".docx":
        return _extract_docx(file_path)
    elif ext == ".xlsx":
        return _extract_xlsx(file_path)
    elif ext == ".hwp":
        return _extract_hwp(file_path)
    elif ext == ".pptx":
        return _extract_pptx(file_path)
    elif ext == ".md":
        return _extract_md(file_path)
    elif ext == ".txt":
        return _extract_txt(file_path)
    else:
        print(f"[skip] 지원하지 않는 파일 형식: {file_path.name}")
        return ""


def extract_content(file_path: Path) -> ExtractionResult:
    """파일에서 텍스트, 표(마크다운), 이미지를 모두 추출한다."""
    ext = file_path.suffix.lower()
    stem = file_path.stem

    if ext == ".pdf":
        return _extract_pdf_content(file_path, stem)
    elif ext == ".docx":
        return _extract_docx_content(file_path, stem)
    elif ext == ".xlsx":
        return _extract_xlsx_content(file_path, stem)
    elif ext == ".pptx":
        return _extract_pptx_content(file_path, stem)
    elif ext == ".md":
        return _extract_md_content(file_path)
    elif ext == ".txt":
        return ExtractionResult(text=_extract_txt(file_path))
    elif ext == ".hwp":
        return ExtractionResult(text=_extract_hwp(file_path))
    else:
        return ExtractionResult(text="")


def extract_from_directory(dir_path: Path) -> list[dict]:
    """폴더 내 지원되는 모든 파일에서 텍스트를 일괄 추출한다."""
    results = []
    if not dir_path.is_dir():
        print(f"[error] 디렉토리가 존재하지 않습니다: {dir_path}")
        return results

    files = sorted(
        f for f in dir_path.iterdir() if f.is_file() and f.suffix.lower() in SUPPORTED_EXTENSIONS
    )
    print(f"[extract] {dir_path}에서 {len(files)}개 파일 발견")

    for file_path in files:
        print(f"  - {file_path.name}")
        try:
            text = extract_text(file_path)
            results.append({
                "filename": file_path.name,
                "extracted_text": text,
            })
        except Exception as e:
            print(f"    [error] 텍스트 추출 실패: {e}")
            results.append({
                "filename": file_path.name,
                "extracted_text": "",
            })

    return results


def extract_pdf_images(file_path: Path, output_dir: Path) -> list[str]:
    """PDF 각 페이지를 PNG 이미지로 변환하여 저장한다. 저장된 파일 경로 목록을 반환."""
    import fitz  # pymupdf

    output_dir.mkdir(parents=True, exist_ok=True)
    stem = file_path.stem
    image_paths = []

    doc = fitz.open(file_path)
    for i, pg in enumerate(doc):
        pix = pg.get_pixmap(dpi=150)
        img_path = output_dir / f"{stem}_p{i + 1}.png"
        pix.save(str(img_path))
        # 프로젝트 루트 기준 상대 경로로 저장
        try:
            rel_path = str(img_path.relative_to(Path(__file__).resolve().parents[2]))
        except ValueError:
            rel_path = str(img_path)
        image_paths.append(rel_path)
    doc.close()

    print(f"    [img] {len(image_paths)}페이지 이미지 저장: {output_dir}")
    return image_paths


def _extract_pdf(file_path: Path) -> str:
    """pdfplumber로 PDF 텍스트를 추출한다."""
    import pdfplumber

    text_parts = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
    return "\n".join(text_parts)


def _extract_docx(file_path: Path) -> str:
    """python-docx로 DOCX 텍스트를 추출한다."""
    from docx import Document

    doc = Document(file_path)
    return "\n".join(para.text for para in doc.paragraphs if para.text.strip())


def _extract_xlsx(file_path: Path) -> str:
    """openpyxl로 XLSX의 모든 시트에서 텍스트를 추출한다."""
    from openpyxl import load_workbook

    wb = load_workbook(file_path, read_only=True, data_only=True)
    text_parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        text_parts.append(f"[시트: {sheet_name}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(cell) for cell in row if cell is not None]
            if cells:
                text_parts.append("\t".join(cells))
    wb.close()
    return "\n".join(text_parts)


def _extract_pptx(file_path: Path) -> str:
    """python-pptx로 PPTX 슬라이드별 텍스트를 추출한다."""
    from pptx import Presentation

    prs = Presentation(file_path)
    text_parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            text_parts.append(f"[슬라이드 {i}]\n" + "\n".join(slide_texts))
    return "\n\n".join(text_parts)


def _extract_md(file_path: Path) -> str:
    """마크다운 파일을 UTF-8 텍스트로 읽는다."""
    return file_path.read_text(encoding="utf-8")


def _extract_txt(file_path: Path) -> str:
    """텍스트 파일을 UTF-8로 읽는다."""
    return file_path.read_text(encoding="utf-8")


def _table_to_markdown(rows: list[list]) -> str:
    """2D 리스트를 마크다운 테이블로 변환한다."""
    if not rows:
        return ""
    # 셀을 문자열로 변환, None → 빈 문자열
    str_rows = [[str(c) if c is not None else "" for c in row] for row in rows]
    if not str_rows:
        return ""
    # 헤더
    header = "| " + " | ".join(str_rows[0]) + " |"
    separator = "| " + " | ".join(["---"] * len(str_rows[0])) + " |"
    body_lines = ["| " + " | ".join(row) + " |" for row in str_rows[1:]]
    return "\n".join([header, separator] + body_lines)


# ── PDF 콘텐츠 추출 (텍스트 + 표 + 이미지) ──

def _extract_pdf_content(file_path: Path, stem: str) -> ExtractionResult:
    """PDF에서 텍스트, 표(마크다운), 임베디드 이미지를 추출한다."""
    import pdfplumber
    import fitz  # pymupdf

    text_parts = []
    images: list[ExtractedImage] = []

    # 텍스트 + 표 (pdfplumber)
    with pdfplumber.open(file_path) as pdf:
        for page_num, page in enumerate(pdf.pages, 1):
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
            # 표 추출
            tables = page.extract_tables()
            for table in tables:
                md_table = _table_to_markdown(table)
                if md_table:
                    text_parts.append(f"\n{md_table}\n")

    # 임베디드 이미지 (PyMuPDF)
    doc = fitz.open(file_path)
    img_counter = 0
    for page_num, fitz_page in enumerate(doc, 1):
        for img_info in fitz_page.get_images(full=True):
            xref = img_info[0]
            try:
                base_image = doc.extract_image(xref)
                if not base_image or not base_image.get("image"):
                    continue
                img_data = base_image["image"]
                img_ext = base_image.get("ext", "png")
                # 작은 이미지 스킵 (아이콘 등)
                if len(img_data) < 2000:
                    continue
                filename = f"{stem}_page{page_num}_img{img_counter}.{img_ext}"
                images.append(ExtractedImage(
                    data=img_data,
                    filename=filename,
                    page_or_slide=page_num,
                    context=f"페이지 {page_num}",
                ))
                img_counter += 1
            except Exception:
                continue
    doc.close()

    return ExtractionResult(text="\n".join(text_parts), images=images)


# ── DOCX 콘텐츠 추출 ──

def _extract_docx_content(file_path: Path, stem: str) -> ExtractionResult:
    """DOCX에서 텍스트, 표, 이미지를 추출한다."""
    from docx import Document
    from docx.opc.constants import RELATIONSHIP_TYPE as RT

    doc = Document(file_path)
    text_parts = []
    images: list[ExtractedImage] = []

    # 텍스트
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)

    # 표
    for table in doc.tables:
        rows = []
        for row in table.rows:
            rows.append([cell.text for cell in row.cells])
        md_table = _table_to_markdown(rows)
        if md_table:
            text_parts.append(f"\n{md_table}\n")

    # 이미지 (relationship에서 추출)
    img_counter = 0
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                img_data = rel.target_part.blob
                if len(img_data) < 2000:
                    continue
                content_type = rel.target_part.content_type
                ext = content_type.split("/")[-1] if "/" in content_type else "png"
                if ext == "jpeg":
                    ext = "jpg"
                filename = f"{stem}_img{img_counter}.{ext}"
                images.append(ExtractedImage(
                    data=img_data,
                    filename=filename,
                    page_or_slide=1,
                    context="문서 이미지",
                ))
                img_counter += 1
            except Exception:
                continue

    return ExtractionResult(text="\n".join(text_parts), images=images)


# ── XLSX 콘텐츠 추출 ──

def _extract_xlsx_content(file_path: Path, stem: str) -> ExtractionResult:
    """XLSX에서 시트별 마크다운 테이블 + 이미지를 추출한다."""
    from openpyxl import load_workbook

    wb = load_workbook(file_path, data_only=True)
    text_parts = []
    images: list[ExtractedImage] = []
    img_counter = 0

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(cells)

        if rows:
            text_parts.append(f"[시트: {sheet_name}]")
            md_table = _table_to_markdown(rows)
            if md_table:
                text_parts.append(md_table)

        # 이미지 추출
        for img in ws._images:
            try:
                img_data = img._data()
                if len(img_data) < 2000:
                    continue
                filename = f"{stem}_{sheet_name}_img{img_counter}.png"
                images.append(ExtractedImage(
                    data=img_data,
                    filename=filename,
                    page_or_slide=1,
                    context=f"시트: {sheet_name}",
                ))
                img_counter += 1
            except Exception:
                continue

    wb.close()
    return ExtractionResult(text="\n".join(text_parts), images=images)


# ── PPTX 콘텐츠 추출 ──

def _extract_pptx_content(file_path: Path, stem: str) -> ExtractionResult:
    """PPTX에서 슬라이드별 텍스트, 표, 이미지를 추출한다."""
    from pptx import Presentation

    prs = Presentation(file_path)
    text_parts = []
    images: list[ExtractedImage] = []
    img_counter = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            # 텍스트
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
            # 표
            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    rows.append([cell.text for cell in row.cells])
                md_table = _table_to_markdown(rows)
                if md_table:
                    slide_texts.append(f"\n{md_table}\n")
            # 이미지
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    img_data = shape.image.blob
                    if len(img_data) < 2000:
                        continue
                    content_type = shape.image.content_type
                    ext = content_type.split("/")[-1] if "/" in content_type else "png"
                    if ext == "jpeg":
                        ext = "jpg"
                    filename = f"{stem}_slide{slide_num}_img{img_counter}.{ext}"
                    slide_context = slide_texts[0] if slide_texts else f"슬라이드 {slide_num}"
                    images.append(ExtractedImage(
                        data=img_data,
                        filename=filename,
                        page_or_slide=slide_num,
                        context=slide_context[:100],
                    ))
                    img_counter += 1
                except Exception:
                    continue

        if slide_texts:
            text_parts.append(f"[슬라이드 {slide_num}]\n" + "\n".join(slide_texts))

    return ExtractionResult(text="\n\n".join(text_parts), images=images)


# ── MD 콘텐츠 추출 ──

def _extract_md_content(file_path: Path) -> ExtractionResult:
    """마크다운 파일에서 텍스트와 참조 이미지를 추출한다."""
    import re

    text = file_path.read_text(encoding="utf-8")
    images: list[ExtractedImage] = []

    # ![alt](path) 패턴의 로컬 이미지 수집
    img_pattern = re.compile(r"!\[([^\]]*)\]\(([^)]+)\)")
    for match in img_pattern.finditer(text):
        alt_text = match.group(1)
        img_ref = match.group(2)
        # URL이 아닌 로컬 경로만 처리
        if img_ref.startswith(("http://", "https://")):
            continue
        img_path = file_path.parent / img_ref
        if img_path.exists():
            try:
                images.append(ExtractedImage(
                    data=img_path.read_bytes(),
                    filename=img_path.name,
                    page_or_slide=1,
                    context=alt_text or "마크다운 이미지",
                ))
            except Exception:
                continue

    return ExtractionResult(text=text, images=images)


def _extract_hwp(file_path: Path) -> str:
    """python-hwp로 HWP 텍스트를 추출한다. 실패 시 olefile 폴백."""
    try:
        from hwp5.hwp5txt import extract_text as hwp5_extract

        return hwp5_extract(str(file_path))
    except Exception:
        pass

    # olefile 폴백
    try:
        import olefile

        if not olefile.isOleFile(str(file_path)):
            print(f"    [warn] OLE 형식이 아닌 HWP: {file_path.name}")
            return ""

        ole = olefile.openole(str(file_path))
        text_parts = []
        for stream in ole.listdir():
            stream_path = "/".join(stream)
            if "BodyText" in stream_path or "PrvText" in stream_path:
                try:
                    data = ole.openstream(stream).read()
                    text = data.decode("utf-16-le", errors="ignore")
                    # NULL 문자 및 제어 문자 제거
                    text = "".join(c for c in text if c.isprintable() or c in "\n\t")
                    if text.strip():
                        text_parts.append(text.strip())
                except Exception:
                    continue
        ole.close()
        return "\n".join(text_parts)
    except Exception as e:
        print(f"    [error] HWP 추출 실패: {e}")
        return ""
